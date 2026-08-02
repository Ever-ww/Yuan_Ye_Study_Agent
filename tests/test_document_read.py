"""统一文档读取工具的格式、分页、边界和 Runtime 装配测试。"""

from __future__ import annotations

import asyncio
import json
import tempfile
import unittest
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject
from pptx import Presentation

from sandbox import WorkspaceLockManager
from tool import ToolContext, default_tools
from tools import (
    DocumentFormatUnsupportedError,
    DocumentReadError,
    DocumentReadResponse,
    ReadFileTool,
)


def _context(root: Path) -> ToolContext:
    return ToolContext(
        project_root=root,
        file_locks=WorkspaceLockManager(root),
    )


def _write_text_pdf(path: Path, pages: list[str], *, title: str = "") -> None:
    writer = PdfWriter()
    font = DictionaryObject({
        NameObject("/Type"): NameObject("/Font"),
        NameObject("/Subtype"): NameObject("/Type1"),
        NameObject("/BaseFont"): NameObject("/Helvetica"),
    })
    font_ref = writer._add_object(font)
    for text in pages:
        page = writer.add_blank_page(width=612, height=792)
        page[NameObject("/Resources")] = DictionaryObject({
            NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref}),
        })
        stream = DecodedStreamObject()
        escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        stream.set_data(f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET".encode("latin-1"))
        page[NameObject("/Contents")] = writer._add_object(stream)
    if title:
        writer.add_metadata({"/Title": title, "/Author": "Researcher"})
    with path.open("wb") as handle:
        writer.write(handle)


class DocumentReadTests(unittest.TestCase):
    def test_pdf_preserves_page_boundaries_and_supports_page_selection(self) -> None:
        async def invoke(root: Path) -> DocumentReadResponse:
            path = root / "paper.pdf"
            _write_text_pdf(
                path,
                ["First page introduction", "Second page experiment results"],
                title="Study Paper",
            )
            raw = await default_tools(root).execute(
                "read_file",
                {"path": "paper.pdf", "start_page": 2, "end_page": 2},
                _context(root),
            )
            return DocumentReadResponse.model_validate_json(raw, strict=True)

        with tempfile.TemporaryDirectory() as value:
            result = asyncio.run(invoke(Path(value)))
        self.assertEqual(result.unit_kind, "page")
        self.assertEqual(result.total_units, 2)
        self.assertEqual((result.selected_start, result.selected_end), (2, 2))
        self.assertEqual(result.title, "Study Paper")
        self.assertEqual(result.metadata["author"], "Researcher")
        self.assertIn("Second page experiment results", result.content)
        self.assertNotIn("First page introduction", result.content)
        self.assertFalse(result.ocr_required)

    def test_image_only_pdf_reports_ocr_requirement(self) -> None:
        async def invoke(root: Path) -> DocumentReadResponse:
            writer = PdfWriter()
            writer.add_blank_page(width=100, height=100)
            with (root / "scan.pdf").open("wb") as handle:
                writer.write(handle)
            raw = await ReadFileTool().run({"path": "scan.pdf"}, _context(root))
            return DocumentReadResponse.model_validate_json(raw, strict=True)

        with tempfile.TemporaryDirectory() as value:
            result = asyncio.run(invoke(Path(value)))
        self.assertTrue(result.ocr_required)
        self.assertTrue(any("OCR" in warning for warning in result.warnings))

    def test_docx_extracts_headings_paragraphs_and_tables(self) -> None:
        async def invoke(root: Path) -> DocumentReadResponse:
            document = Document()
            document.core_properties.title = "项目说明"
            document.add_heading("研究目标", level=1)
            document.add_paragraph("分析长上下文 Agent。")
            table = document.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "指标"
            table.cell(0, 1).text = "值"
            table.cell(1, 0).text = "准确率"
            table.cell(1, 1).text = "95%"
            document.save(root / "study.docx")
            raw = await ReadFileTool().run({"path": "study.docx"}, _context(root))
            return DocumentReadResponse.model_validate_json(raw, strict=True)

        with tempfile.TemporaryDirectory() as value:
            result = asyncio.run(invoke(Path(value)))
        self.assertEqual(result.title, "项目说明")
        self.assertIn("# 研究目标", result.content)
        self.assertIn("分析长上下文 Agent", result.content)
        self.assertIn("| 指标 | 值 |", result.content)
        self.assertIn("| 准确率 | 95% |", result.content)

    def test_pptx_and_xlsx_use_slides_and_worksheets_as_units(self) -> None:
        async def invoke(root: Path) -> tuple[DocumentReadResponse, DocumentReadResponse, DocumentReadResponse]:
            deck = Presentation()
            first = deck.slides.add_slide(deck.slide_layouts[1])
            first.shapes.title.text = "第一页"
            first.placeholders[1].text = "背景"
            second = deck.slides.add_slide(deck.slide_layouts[1])
            second.shapes.title.text = "第二页"
            second.placeholders[1].text = "结论"
            deck.save(root / "talk.pptx")

            workbook = Workbook()
            workbook.active.title = "Data"
            workbook.active.append(["metric", "value"])
            workbook.active.append(["accuracy", 0.95])
            workbook.create_sheet("Notes")["A1"] = "reviewed"
            workbook.save(root / "results.xlsx")

            ppt_raw = await ReadFileTool().run(
                {"path": "talk.pptx", "start_page": 2, "end_page": 2},
                _context(root),
            )
            xlsx_raw = await ReadFileTool().run(
                {"path": "results.xlsx", "start_page": 1, "end_page": 1},
                _context(root),
            )
            return (
                DocumentReadResponse.model_validate_json(ppt_raw, strict=True),
                DocumentReadResponse.model_validate_json(xlsx_raw, strict=True),
            )

        with tempfile.TemporaryDirectory() as value:
            slides, sheets = asyncio.run(invoke(Path(value)))
        self.assertEqual(slides.unit_kind, "slide")
        self.assertIn("第二页", slides.content)
        self.assertNotIn("第一页", slides.content)
        self.assertEqual(sheets.unit_kind, "worksheet")
        self.assertIn("metric\tvalue", sheets.content)
        self.assertNotIn("reviewed", sheets.content)

    def test_html_notebook_truncation_and_offset_are_controlled(self) -> None:
        async def invoke(root: Path) -> tuple[DocumentReadResponse, DocumentReadResponse]:
            (root / "page.html").write_text(
                "<html><body><h1>标题</h1><script>secret()</script><p>" + "A" * 1600 + "</p></body></html>",
                encoding="utf-8",
            )
            (root / "notes.ipynb").write_text(json.dumps({
                "cells": [
                    {"cell_type": "markdown", "source": ["# 说明"]},
                    {"cell_type": "code", "source": ["print('ok')"]},
                ],
            }), encoding="utf-8")
            first_raw = await ReadFileTool().run(
                {"path": "page.html", "max_chars": 1000},
                _context(root),
            )
            continued_raw = await ReadFileTool().run(
                {"path": "page.html", "offset_chars": 1000, "max_chars": 1000},
                _context(root),
            )
            notebook_raw = await ReadFileTool().run({"path": "notes.ipynb"}, _context(root))
            return (
                DocumentReadResponse.model_validate_json(first_raw, strict=True),
                DocumentReadResponse.model_validate_json(continued_raw, strict=True),
                DocumentReadResponse.model_validate_json(notebook_raw, strict=True),
            )

        with tempfile.TemporaryDirectory() as value:
            html, continued, notebook = asyncio.run(invoke(Path(value)))
        self.assertTrue(html.truncated)
        self.assertEqual(html.next_offset_chars, 1000)
        self.assertEqual(continued.offset_chars, 1000)
        self.assertFalse(continued.truncated)
        self.assertGreater(len(continued.content), 500)
        self.assertNotIn("secret", html.content)
        self.assertTrue(html.untrusted_document_content)
        self.assertIn("Cell 1 (markdown)", notebook.content)
        self.assertIn("print('ok')", notebook.content)

    def test_legacy_unknown_corrupt_and_workspace_escape_are_rejected(self) -> None:
        async def invoke(root: Path) -> None:
            tool = ReadFileTool()
            context = _context(root)
            (root / "old.doc").write_bytes(b"legacy")
            (root / "image.png").write_bytes(b"png")
            (root / "broken.docx").write_bytes(b"not-a-zip")
            with self.assertRaises(DocumentFormatUnsupportedError):
                await tool.run({"path": "old.doc"}, context)
            with self.assertRaises(DocumentFormatUnsupportedError):
                await tool.run({"path": "image.png"}, context)
            with self.assertRaises(DocumentReadError):
                await tool.run({"path": "broken.docx"}, context)
            with self.assertRaises(PermissionError):
                await tool.run({"path": "../outside.pdf"}, context)

        with tempfile.TemporaryDirectory() as value:
            asyncio.run(invoke(Path(value)))

    def test_default_registry_exposes_only_unified_read_file(self) -> None:
        with tempfile.TemporaryDirectory() as value:
            root = Path(value)
            registry = default_tools(root)
            self.assertIn("read_file", registry.names())
            self.assertNotIn("read_document", registry.names())
            schema = next(item for item in registry.schemas() if item["name"] == "read_file")
            self.assertEqual(schema["parameters"]["required"], ["path"])


if __name__ == "__main__":
    unittest.main()
